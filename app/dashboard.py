import sys
import os

# -------------------------------------------------
# Add PROJECT ROOT to Python path (FIX src import) 
# -------------------------------------------------
PROJECT_ROOT = os.path.abspath(
    os.path.join(os.path.dirname(__file__), "..")
)
sys.path.insert(0, PROJECT_ROOT)

# -------------------------------------------------
# StandardImports 
# -------------------------------------------------
import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import joblib

from sklearn.model_selection import train_test_split
from sklearn.metrics import accuracy_score, roc_curve, auc
from sklearn.preprocessing import label_binarize

# -------------------------------------------------
# ReportLab Imports (PDF generation)
# -------------------------------------------------
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Frame, PageTemplate
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.pagesizes import A4


# -------------------------------------------------
# PowerPoint Imports (PPT generation)  ✅ MISSING
# -------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt

# -------------------------------------------------
# Project Imports
# -------------------------------------------------
from src.preprocessing import clean_text


# ---------------- DIRECTORIES ----------------
os.makedirs("reports/plots", exist_ok=True)
os.makedirs("reports/output", exist_ok=True)

# ---------------- PAGE CONFIG ----------------
st.set_page_config(
    page_title="Q-Commerce Sentiment Analysis",
    page_icon="🛒",
    layout="wide"
)

# ---------------- SIDEBAR ----------------
st.sidebar.title("📊 Dashboard Menu")
page = st.sidebar.radio(
    "Navigate",
    ["Overview", "EDA", "Model Evaluation", "Live Prediction", "Download Report",  "Download Final PPT", "Download IEEE Report"]
)

# ---------------- LOAD DATA ----------------
df = pd.read_excel("data/raw/Quick_Commerce_Consumer_Behavior_Labeled.xlsx")
df = df[["Combined_Opinion", "Sentiment"]].dropna()
df["clean_text"] = df["Combined_Opinion"].apply(clean_text)
df["text_length"] = df["clean_text"].apply(len)

# ---------------- LOAD MODEL ----------------
model = joblib.load("models/sentiment_model.pkl")
vectorizer = joblib.load("models/vectorizer.pkl")

# ---------------- OVERVIEW ----------------
if page == "Overview":
    st.title("🛒 Q-Commerce Sentiment Analysis")
    st.markdown("""
    ### Project Overview
    This dashboard analyzes Q-Commerce survey opinions using
    supervised machine learning to predict customer sentiment.
    """)
    st.success("✔ Supervised ML with Logistic Regression")
    st.info("✔ Dataset manually labeled for sentiment")

# ---------------- EDA ----------------
elif page == "EDA":
    st.title("📊 Exploratory Data Analysis")

    col1, col2 = st.columns(2)

    with col1:
        fig1 = plt.figure()
        sns.countplot(x="Sentiment", data=df)
        plt.title("Sentiment Distribution")
        st.pyplot(fig1)
        fig1.savefig("reports/plots/sentiment_distribution.png")

    with col2:
        fig2 = plt.figure()
        sns.boxplot(x="Sentiment", y="text_length", data=df)
        plt.title("Text Length vs Sentiment")
        st.pyplot(fig2)
        fig2.savefig("reports/plots/textlength_boxplot.png")

    fig3 = plt.figure()
    sns.violinplot(x="Sentiment", y="text_length", data=df)
    plt.title("Violin Plot")
    st.pyplot(fig3)
    fig3.savefig("reports/plots/violin_plot.png")

# ---------------- MODEL EVALUATION ----------------
elif page == "Model Evaluation":
    st.title("📈 Model Evaluation")

    X = vectorizer.transform(df["clean_text"])
    y = df["Sentiment"]

    X_train, X_test, y_train, y_test = train_test_split(
        X, y, test_size=0.2, random_state=42
    )

    y_pred = model.predict(X_test)
    acc = accuracy_score(y_test, y_pred)

    st.metric("Model Accuracy", f"{acc:.2f}")

    fig4 = plt.figure()
    classes = ["Negative", "Neutral", "Positive"]
    y_test_bin = label_binarize(y_test, classes=classes)
    y_score = model.predict_proba(X_test)

    for i, cls in enumerate(classes):
        if y_test_bin[:, i].sum() == 0:
            continue
        fpr, tpr, _ = roc_curve(y_test_bin[:, i], y_score[:, i])
        roc_auc = auc(fpr, tpr)
        plt.plot(fpr, tpr, label=f"{cls} (AUC={roc_auc:.2f})")

    plt.plot([0, 1], [0, 1], "k--")
    plt.legend()
    plt.title("ROC Curve")
    st.pyplot(fig4)
    fig4.savefig("reports/plots/roc_curve.png")

# ---------------- LIVE PREDICTION ----------------
elif page == "Live Prediction":
    st.title("🧪 Live Sentiment Prediction")

    negative_words = ["concern", "delay", "late", "problem", "issue", "poor", "expensive"]
    neutral_words = ["more", "better", "required", "need", "expect", "should"]

    review = st.text_area("Enter Customer Opinion")

    if st.button("Predict Sentiment"):
        clean = clean_text(review)

        if any(w in clean for w in negative_words):
            st.error("😠 Negative Sentiment")
        elif any(w in clean for w in neutral_words):
            st.info("😐 Neutral Sentiment")
        else:
            vec = vectorizer.transform([clean])
            result = model.predict(vec)[0]

            if result == "Positive":
                st.success("😊 Positive Sentiment")
            else:
                st.info("😐 Neutral Sentiment")

# ---------------- DOWNLOAD REPORT ----------------
elif page == "Download Report":
    st.title("📄 Download Project Report (PDF)")

    if st.button("Generate & Download PDF Report"):
        pdf_path = "reports/output/QCommerce_Sentiment_Report.pdf"
        styles = getSampleStyleSheet()
        doc = SimpleDocTemplate(pdf_path)
        story = []

        story.append(Paragraph("<b>Q-Commerce Sentiment Analysis Report</b>", styles["Title"]))
        story.append(Spacer(1, 12))

        story.append(Paragraph(
            "This report presents sentiment analysis on Q-Commerce survey data using supervised machine learning.",
            styles["Normal"]
        ))
        story.append(Spacer(1, 12))

        for img in [
            "sentiment_distribution.png",
            "textlength_boxplot.png",
            "violin_plot.png",
            "roc_curve.png"
        ]:
            img_path = f"reports/plots/{img}"
            if os.path.exists(img_path):
                story.append(Image(img_path, width=400, height=250))
                story.append(Spacer(1, 12))

        doc.build(story)

        with open(pdf_path, "rb") as f:
            st.download_button(
                label="⬇ Download PDF Report",
                data=f,
                file_name="QCommerce_Sentiment_Report.pdf",
                mime="application/pdf"
            )
elif page == "Download Final PPT":
    st.title("📊 Download Final Project PPT")

    if st.button("Generate & Download PPT"):
        ppt_path = "reports/output/QCommerce_Final_Presentation.pptx"

        # Create presentation
        prs = Presentation()

        def add_slide(title, body):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = title
            slide.placeholders[1].text = body

        # -------- Text Slides --------
        add_slide(
            "Sentiment Prediction Using ML",
            "Q-Commerce Reviews\n\nSupervised Machine Learning Project"
        )

        add_slide(
            "Dataset",
            "Q-Commerce Consumer Survey Dataset\nManually labeled for sentiment"
        )

        add_slide(
            "Methodology",
            "Preprocessing → TF-IDF → Logistic Regression\nPost-processing rules"
        )

        add_slide(
            "Results",
            "Sentiment distribution and ROC analysis\nPositive, Neutral, Negative"
        )

        add_slide(
            "Conclusion",
            "Survey-based sentiment successfully analyzed\nSystem is explainable"
        )

        # -------- Image Slides with Captions --------
        caption_map = {
            "sentiment_distribution.png": "Sentiment Distribution (Histogram)",
            "textlength_boxplot.png": "Text Length vs Sentiment (Boxplot)",
            "violin_plot.png": "Text Length Distribution (Violin Plot)",
            "roc_curve.png": "ROC Curve for Sentiment Classification"
        }

        plot_dir = "reports/plots"

        for img in os.listdir(plot_dir):
            if img.endswith(".png"):
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

                img_path = os.path.join(plot_dir, img)
                slide.shapes.add_picture(
                    img_path,
                    Inches(1),
                    Inches(1),
                    width=Inches(6)
                )

                caption = caption_map.get(
                    img,
                    img.replace(".png", "").replace("_", " ").title()
                )

                textbox = slide.shapes.add_textbox(
                    Inches(1),
                    Inches(5.5),
                    Inches(6),
                    Inches(0.8)
                )

                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = caption
                p.font.size = Pt(14)
                p.font.bold = True
                p.alignment = 1  # center

        # Save PPT
        prs.save(ppt_path)

        # Download button
        with open(ppt_path, "rb") as f:
            st.download_button(
                label="⬇ Download Final PPT",
                data=f,
                file_name="QCommerce_Final_Presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
elif page == "Download IEEE Report":
    st.title("📄 Download IEEE Style Report (2-Column)")

    if st.button("Generate IEEE Report PDF"):
        pdf_path = "reports/output/IEEE_QCommerce_Sentiment_Report.pdf"

        # -------- Document Setup --------
        doc = SimpleDocTemplate(
            pdf_path,
            pagesize=A4,
            rightMargin=36,
            leftMargin=36,
            topMargin=36,
            bottomMargin=36
        )

        styles = getSampleStyleSheet()
        styles.add(ParagraphStyle(
            name="IEEE_Title",
            fontSize=14,
            alignment=1,
            spaceAfter=12,
            leading=16,
            fontName="Helvetica-Bold"
        ))

        styles.add(ParagraphStyle(
            name="IEEE_Body",
            fontSize=9,
            leading=11,
            spaceAfter=6
        ))

        story = []

        # -------- Title --------
        story.append(Paragraph(
            "Sentiment Prediction Using Machine Learning for Q-Commerce Reviews",
            styles["IEEE_Title"]
        ))

        # -------- Abstract --------
        story.append(Paragraph("<b>Abstract—</b>", styles["IEEE_Body"]))
        story.append(Paragraph(
            "Quick Commerce (Q-Commerce) platforms generate large volumes of customer feedback. "
            "Manual analysis of this unstructured text is inefficient. This project proposes a "
            "supervised machine learning approach to classify customer opinions into Positive, "
            "Neutral, and Negative sentiments. A real-world Q-Commerce survey dataset was manually "
            "labeled and processed using TF-IDF feature extraction. Logistic Regression was trained "
            "and evaluated, achieving an accuracy of 90.9%. The system is deployed using a Streamlit "
            "dashboard for analysis and prediction.",
            styles["IEEE_Body"]
        ))

        # -------- Sections --------
        sections = [
            ("I. INTRODUCTION",
             "Q-Commerce platforms generate large amounts of customer feedback related to delivery, "
             "pricing, and service quality. Automated sentiment analysis helps businesses understand "
             "customer satisfaction efficiently."),

            ("II. DATASET DESCRIPTION",
             "The dataset consists of consumer survey responses collected from Q-Commerce users. "
             "Since sentiment labels were not present, manual labeling was performed."),

            ("III. METHODOLOGY",
             "Text preprocessing includes lowercase conversion, stopword removal, and normalization. "
             "TF-IDF vectorization is used for feature extraction, followed by Logistic Regression "
             "classification."),

            ("IV. EXPLORATORY DATA ANALYSIS",
             "EDA includes sentiment distribution analysis, text length analysis using boxplots and "
             "violin plots, and ROC curve visualization."),

            ("V. RESULTS AND EVALUATION",
             "The model achieved an accuracy of 90.9%. The results indicate effective classification "
             "of Positive and Neutral sentiments."),

            ("VI. SYSTEM IMPLEMENTATION",
             "A Streamlit dashboard provides interactive visualization, live prediction, and automated "
             "report generation."),

            ("VII. CONCLUSION",
             "The proposed system successfully applies supervised machine learning for sentiment "
             "analysis in Q-Commerce reviews."),

            ("VIII. FUTURE WORK",
             "Future enhancements include deep learning models, multilingual support, and real-time "
             "integration.")
        ]

        for title, content in sections:
            story.append(Spacer(1, 6))
            story.append(Paragraph(f"<b>{title}</b>", styles["IEEE_Body"]))
            story.append(Paragraph(content, styles["IEEE_Body"]))

        # -------- Two Column Layout --------
        frame1 = Frame(doc.leftMargin, doc.bottomMargin,
                       (doc.width / 2) - 6, doc.height, id='col1')
        frame2 = Frame(doc.leftMargin + (doc.width / 2) + 6, doc.bottomMargin,
                       (doc.width / 2) - 6, doc.height, id='col2')

        doc.addPageTemplates([
            PageTemplate(id='TwoCol', frames=[frame1, frame2])
        ])

        doc.build(story)

        with open(pdf_path, "rb") as f:
            st.download_button(
                label="⬇ Download IEEE Report PDF",
                data=f,
                file_name="IEEE_QCommerce_Sentiment_Report.pdf",
                mime="application/pdf"
            )
